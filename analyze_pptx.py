from pptx import Presentation
from pptx.util import Inches, Emu

prs = Presentation(r'c:\Users\Harsh\OneDrive\Desktop\Agromod_Prototype_Submission.pptx')
print(f'Slide width: {prs.slide_width} emu = {prs.slide_width / 914400:.2f} inches')
print(f'Slide height: {prs.slide_height} emu = {prs.slide_height / 914400:.2f} inches')
print(f'Total slides: {len(prs.slides)}')
print()
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for j, shape in enumerate(slide.shapes):
        print(f'  Shape {j}: type={shape.shape_type}, name={shape.name}')
        print(f'    pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if hasattr(shape, 'image'):
            try:
                print(f'    IMAGE: content_type={shape.image.content_type}, size={len(shape.image.blob)} bytes')
            except:
                pass
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text[:120] if p.text else ''
                if txt:
                    font_info = ''
                    if p.runs:
                        r = p.runs[0]
                        try:
                            color = r.font.color.rgb if r.font.color and r.font.color.rgb else 'inherit'
                        except:
                            color = 'theme'
                        font_info = f' [font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={color}]'
                    print(f'    TEXT: "{txt}"{font_info}')
    print()
